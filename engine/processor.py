#!/usr/bin/env python3
"""
LawPro Fast Converter - File Processor
=======================================
하이브리드 문서 처리 라우터

Routing Logic:
- xlsx, docx, pptx, hwpx, digital PDF → Local Libraries (비용 0원, 초고속)
- image PDF → Upstage Document Parse API (고정밀 OCR)
"""

import os
import io
import re
import time
import base64
import zipfile
from typing import Dict, Any, Optional, Tuple
from xml.etree import ElementTree as ET

import requests

# Clean HTML 및 마크다운 변환용
try:
    from cleaner import ContentCleaner
    HAS_CLEANER = True
except ImportError:
    HAS_CLEANER = False

# 크레딧 관리
try:
    from credit_manager import get_credit_manager
    HAS_CREDIT = True
except ImportError:
    HAS_CREDIT = False

# 관리자 설정 (Upstage API 키)
try:
    from admin_config import get_admin_config
    HAS_ADMIN_CONFIG = True
except ImportError:
    HAS_ADMIN_CONFIG = False

# Rate Limiter (적응형 Rate Limit 학습)
try:
    from rate_limiter import get_rate_limiter
    HAS_RATE_LIMITER = True
except ImportError:
    HAS_RATE_LIMITER = False


class FileProcessor:
    """하이브리드 문서 처리기"""

    # Upstage API 설정
    UPSTAGE_API_URL = "https://api.upstage.ai/v1/document-ai/document-parse"
    MAX_PDF_PAGES_PER_REQUEST = 10  # Upstage 권장: 최대 10페이지
    MAX_FILE_SIZE_MB = 50  # 최대 파일 크기 (MB)
    # 429 에러 발생 시 재시도 대기 시간 (초) - 충분히 길게
    RATE_LIMIT_RETRY_DELAYS = [10, 30, 60, 120, 180]  # 10초, 30초, 1분, 2분, 3분

    def __init__(self, output_folder: str,
                 generate_clean_html: bool = True,
                 generate_markdown: bool = True,
                 check_credits: bool = True,
                 api_key: str = None):  # api_key는 하위호환을 위해 유지 (사용되지 않음)
        """
        FileProcessor 초기화

        Args:
            output_folder: 출력 폴더 경로
            generate_clean_html: Clean HTML 생성 여부
            generate_markdown: Markdown 생성 여부
            check_credits: 크레딧 확인 여부
            api_key: (deprecated) 하위호환용 - 관리자 설정에서 자동 로드됨
        """
        self.output_folder = output_folder
        self.generate_clean_html = generate_clean_html
        self.generate_markdown = generate_markdown
        self.check_credits = check_credits

        # 관리자 설정에서 Upstage API 키 로드
        if HAS_ADMIN_CONFIG:
            admin_config = get_admin_config()
            self.api_key = admin_config.upstage_api_key
        else:
            self.api_key = api_key or ""

        # Cleaner 초기화
        if HAS_CLEANER and (generate_clean_html or generate_markdown):
            self.cleaner = ContentCleaner()
        else:
            self.cleaner = None

        # Credit Manager 초기화
        if HAS_CREDIT and check_credits:
            self.credit_manager = get_credit_manager()
        else:
            self.credit_manager = None

        # Rate Limiter 초기화 (적응형 Rate Limit 학습)
        if HAS_RATE_LIMITER:
            self.rate_limiter = get_rate_limiter()
        else:
            self.rate_limiter = None

    def process(self, file_path: str) -> Dict[str, Any]:
        """
        파일 처리 메인 함수

        Returns:
            Dict with keys: status, file, method, time, error (optional)

        Output Structure:
            {output_folder}/{filename}/
                ├── view.html       (원본 서식 - 열람용)
                ├── clean_ai.html   (AI 학습용 - JS/CSS 제거)
                └── content.md      (마크다운 - 노트앱 호환)
        """
        filename = os.path.basename(file_path)
        doc_name = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1].lower()

        # 파일별 출력 폴더 생성
        doc_folder = os.path.join(self.output_folder, doc_name)
        os.makedirs(doc_folder, exist_ok=True)

        start_time = time.time()
        method = "Local"
        outputs = []

        try:
            content = ""

            # 확장자별 라우팅
            if ext in ('.xlsx', '.xls'):
                content = self._convert_excel(file_path)

            elif ext in ('.docx', '.doc'):
                content = self._convert_word(file_path)

            elif ext in ('.pptx', '.ppt'):
                content = self._convert_powerpoint(file_path)

            elif ext in ('.hwpx', '.hwp'):
                content = self._convert_hwp(file_path)

            elif ext == '.pdf':
                is_digital, text_ratio = self._analyze_pdf(file_path)
                if is_digital:
                    content = self._convert_digital_pdf(file_path)
                else:
                    method = "Upstage API"
                    content = self._convert_image_pdf_upstage(file_path)

            else:
                raise ValueError(f"지원하지 않는 파일 형식: {ext}")

            # === 3-Way Output 저장 ===

            # 1. View HTML (원본 서식용)
            view_path = os.path.join(doc_folder, "view.html")
            self._save_html(view_path, content, filename)
            outputs.append("view.html")

            # 2. Clean HTML (AI 학습용)
            if self.generate_clean_html and self.cleaner:
                clean_path = os.path.join(doc_folder, "clean_ai.html")
                clean_html = self.cleaner.make_clean_html_for_ai(content)
                with open(clean_path, 'w', encoding='utf-8') as f:
                    f.write(clean_html)
                outputs.append("clean_ai.html")

            # 3. Markdown (노트앱 호환용)
            if self.generate_markdown and self.cleaner:
                md_path = os.path.join(doc_folder, "content.md")
                # Clean HTML 사용 (없으면 원본에서 변환)
                if self.generate_clean_html:
                    markdown = self.cleaner.convert_to_markdown(clean_html)
                else:
                    clean_html = self.cleaner.make_clean_html_for_ai(content)
                    markdown = self.cleaner.convert_to_markdown(clean_html)
                with open(md_path, 'w', encoding='utf-8') as f:
                    f.write(markdown)
                outputs.append("content.md")

            elapsed = round(time.time() - start_time, 2)
            return {
                "status": "success",
                "file": filename,
                "method": method,
                "time": elapsed,
                "output": doc_folder,
                "outputs": outputs
            }

        except Exception as e:
            return {
                "status": "fail",
                "file": filename,
                "method": method,
                "error": str(e)
            }

    # ============================================================
    # Excel 처리
    # ============================================================
    def _convert_excel(self, file_path: str) -> str:
        """Excel을 HTML 테이블로 변환 (서식 완벽 보존)"""
        try:
            import pandas as pd
            from openpyxl import load_workbook
            from openpyxl.utils import get_column_letter
        except ImportError:
            return self._fallback_excel(file_path)

        wb = load_workbook(file_path, data_only=True)
        html_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            html_parts.append(f'<div class="sheet" data-sheet="{sheet_name}">')
            html_parts.append(f'<h2 class="sheet-title">{sheet_name}</h2>')

            # 병합 셀 정보 수집
            merged_ranges = {}
            for merged_range in ws.merged_cells.ranges:
                min_col, min_row, max_col, max_row = merged_range.bounds
                merged_ranges[(min_row, min_col)] = (
                    max_row - min_row + 1,
                    max_col - min_col + 1
                )

            html_parts.append('<table class="excel-table">')

            for row_idx, row in enumerate(ws.iter_rows(), start=1):
                html_parts.append('<tr>')
                for col_idx, cell in enumerate(row, start=1):
                    # 병합된 셀 처리
                    if (row_idx, col_idx) in merged_ranges:
                        rowspan, colspan = merged_ranges[(row_idx, col_idx)]
                        span_attrs = f' rowspan="{rowspan}" colspan="{colspan}"'
                    else:
                        # 병합 범위 내 다른 셀은 스킵
                        skip = False
                        for (mr, mc), (rs, cs) in merged_ranges.items():
                            if mr <= row_idx < mr + rs and mc <= col_idx < mc + cs:
                                if (row_idx, col_idx) != (mr, mc):
                                    skip = True
                                    break
                        if skip:
                            continue
                        span_attrs = ""

                    # 셀 스타일 추출
                    style = self._extract_cell_style(cell)
                    value = cell.value if cell.value is not None else ""

                    html_parts.append(f'<td{span_attrs} style="{style}">{value}</td>')

                html_parts.append('</tr>')

            html_parts.append('</table></div>')

        return '\n'.join(html_parts)

    def _extract_cell_style(self, cell) -> str:
        """셀 스타일 추출"""
        styles = []

        try:
            # 배경색
            if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb:
                if cell.fill.fgColor.rgb != '00000000':
                    color = cell.fill.fgColor.rgb
                    if len(color) == 8:
                        color = color[2:]  # ARGB -> RGB
                    styles.append(f'background-color: #{color}')

            # 글꼴
            if cell.font:
                if cell.font.bold:
                    styles.append('font-weight: bold')
                if cell.font.italic:
                    styles.append('font-style: italic')
                if cell.font.size:
                    styles.append(f'font-size: {cell.font.size}pt')
                if cell.font.color and cell.font.color.rgb:
                    color = cell.font.color.rgb
                    if len(color) == 8:
                        color = color[2:]
                    if color != '000000':
                        styles.append(f'color: #{color}')

            # 정렬
            if cell.alignment:
                if cell.alignment.horizontal:
                    styles.append(f'text-align: {cell.alignment.horizontal}')
                if cell.alignment.vertical:
                    styles.append(f'vertical-align: {cell.alignment.vertical}')

            # 테두리
            if cell.border:
                border_style = self._get_border_style(cell.border)
                if border_style:
                    styles.append(border_style)

        except Exception:
            pass

        return '; '.join(styles)

    def _get_border_style(self, border) -> str:
        """테두리 스타일 변환"""
        borders = []

        sides = [
            ('top', border.top),
            ('right', border.right),
            ('bottom', border.bottom),
            ('left', border.left)
        ]

        for side_name, side in sides:
            if side and side.style and side.style != 'none':
                width = '1px' if side.style == 'thin' else '2px' if side.style in ('medium', 'thick') else '1px'
                line_style = 'dashed' if side.style == 'dashed' else 'dotted' if side.style == 'dotted' else 'solid'
                borders.append(f'border-{side_name}: {width} {line_style} #000')

        return '; '.join(borders)

    def _fallback_excel(self, file_path: str) -> str:
        """Pandas 폴백"""
        import pandas as pd
        dfs = pd.read_excel(file_path, sheet_name=None)
        html_parts = []
        for sheet_name, df in dfs.items():
            html_parts.append(f'<h2>{sheet_name}</h2>')
            html_parts.append(df.to_html(index=False, border=1, classes='excel-table'))
        return '\n'.join(html_parts)

    # ============================================================
    # Word 처리
    # ============================================================
    def _convert_word(self, file_path: str) -> str:
        """Word 문서를 HTML로 변환"""
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
        except ImportError:
            return self._convert_word_basic(file_path)

        doc = Document(file_path)
        html_parts = []

        for element in doc.element.body:
            if element.tag.endswith('p'):
                # 단락 처리
                for para in doc.paragraphs:
                    if para._element == element:
                        html_parts.append(self._convert_paragraph(para))
                        break

            elif element.tag.endswith('tbl'):
                # 테이블 처리
                for table in doc.tables:
                    if table._tbl == element:
                        html_parts.append(self._convert_table(table))
                        break

        return '\n'.join(html_parts)

    def _convert_paragraph(self, para) -> str:
        """단락을 HTML로 변환"""
        if not para.text.strip():
            return '<p>&nbsp;</p>'

        # 스타일 감지
        style_name = para.style.name if para.style else ""
        if 'Heading 1' in style_name:
            tag = 'h1'
        elif 'Heading 2' in style_name:
            tag = 'h2'
        elif 'Heading 3' in style_name:
            tag = 'h3'
        else:
            tag = 'p'

        # 정렬
        align = ""
        if para.alignment:
            align_map = {0: 'left', 1: 'center', 2: 'right', 3: 'justify'}
            align_val = align_map.get(para.alignment, '')
            if align_val:
                align = f' style="text-align: {align_val}"'

        # 런 단위 서식
        content = ""
        for run in para.runs:
            text = run.text
            if run.bold:
                text = f'<strong>{text}</strong>'
            if run.italic:
                text = f'<em>{text}</em>'
            if run.underline:
                text = f'<u>{text}</u>'
            content += text

        return f'<{tag}{align}>{content}</{tag}>'

    def _convert_table(self, table) -> str:
        """테이블을 HTML로 변환 (테두리 유무 감지)"""
        # 테두리 스타일 감지
        has_border = self._detect_table_border(table)
        border_class = 'bordered' if has_border else 'borderless'

        html = f'<table class="word-table {border_class}">'

        for row in table.rows:
            html += '<tr>'
            for cell in row.cells:
                # 셀 병합 처리
                cell_html = '<td'

                # 수직 병합
                if hasattr(cell, '_tc'):
                    vmerge = cell._tc.xpath('.//w:vMerge')
                    if vmerge:
                        # 병합 시작 셀이 아니면 스킵
                        if vmerge[0].get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val') is None:
                            continue

                # 셀 배경색
                style = self._get_cell_shading(cell)
                if style:
                    cell_html += f' style="{style}"'

                cell_html += '>'
                cell_html += '<br>'.join([p.text or '' for p in cell.paragraphs])
                cell_html += '</td>'
                html += cell_html
            html += '</tr>'

        html += '</table>'
        return html

    def _detect_table_border(self, table) -> bool:
        """테이블 테두리 유무 감지"""
        try:
            tbl_pr = table._tbl.xpath('.//w:tblBorders')
            if tbl_pr:
                for border in tbl_pr[0]:
                    val = border.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')
                    if val and val not in ('nil', 'none'):
                        return True
            return False
        except:
            return True  # 기본값은 테두리 있음

    def _get_cell_shading(self, cell) -> str:
        """셀 음영 추출"""
        try:
            shading = cell._tc.xpath('.//w:shd')
            if shading:
                fill = shading[0].get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
                if fill and fill.lower() not in ('auto', 'ffffff'):
                    return f'background-color: #{fill}'
        except:
            pass
        return ""

    def _convert_word_basic(self, file_path: str) -> str:
        """기본 Word 변환 (python-docx 없을 때)"""
        # DOCX는 ZIP 파일 구조
        if file_path.endswith('.docx'):
            with zipfile.ZipFile(file_path) as zf:
                if 'word/document.xml' in zf.namelist():
                    xml_content = zf.read('word/document.xml')
                    tree = ET.fromstring(xml_content)

                    # 텍스트만 추출
                    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                    paragraphs = tree.findall('.//w:p', ns)

                    html = ""
                    for para in paragraphs:
                        texts = para.findall('.//w:t', ns)
                        text = ''.join([t.text or '' for t in texts])
                        if text.strip():
                            html += f'<p>{text}</p>\n'

                    return html

        return "<p>문서를 읽을 수 없습니다. python-docx 라이브러리를 설치해주세요.</p>"

    # ============================================================
    # PowerPoint 처리
    # ============================================================
    def _convert_powerpoint(self, file_path: str) -> str:
        """PowerPoint를 HTML로 변환"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return self._convert_pptx_basic(file_path)

        prs = Presentation(file_path)
        html_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            html_parts.append(f'<div class="slide" data-slide="{slide_num}">')
            html_parts.append(f'<h2 class="slide-number">Slide {slide_num}</h2>')

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = ''.join([run.text for run in para.runs])
                        if text.strip():
                            # 제목 감지
                            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                                html_parts.append(f'<h3 class="slide-title">{text}</h3>')
                            else:
                                html_parts.append(f'<p>{text}</p>')

                elif shape.has_table:
                    html_parts.append(self._convert_pptx_table(shape.table))

            html_parts.append('</div><hr class="slide-divider"/>')

        return '\n'.join(html_parts)

    def _convert_pptx_table(self, table) -> str:
        """PPTX 테이블 변환"""
        html = '<table class="pptx-table">'
        for row in table.rows:
            html += '<tr>'
            for cell in row.cells:
                text = cell.text if cell.text else ''
                html += f'<td>{text}</td>'
            html += '</tr>'
        html += '</table>'
        return html

    def _convert_pptx_basic(self, file_path: str) -> str:
        """기본 PPTX 변환"""
        html_parts = []

        with zipfile.ZipFile(file_path) as zf:
            # 슬라이드 파일 찾기
            slide_files = sorted([f for f in zf.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])

            for i, slide_file in enumerate(slide_files, 1):
                xml_content = zf.read(slide_file)
                tree = ET.fromstring(xml_content)

                # 모든 텍스트 추출
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                texts = tree.findall('.//a:t', ns)

                html_parts.append(f'<div class="slide"><h2>Slide {i}</h2>')
                for t in texts:
                    if t.text and t.text.strip():
                        html_parts.append(f'<p>{t.text}</p>')
                html_parts.append('</div><hr/>')

        return '\n'.join(html_parts)

    # ============================================================
    # HWP/HWPX 처리
    # ============================================================
    def _convert_hwp(self, file_path: str) -> str:
        """한글 문서를 HTML로 변환"""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.hwpx':
            return self._convert_hwpx(file_path)
        else:
            # .hwp는 바이너리 형식으로 복잡함
            # Upstage API 사용 권장
            if self.api_key:
                return self._convert_image_pdf_upstage(file_path)
            return "<p>HWP 파일은 Upstage API 키가 필요합니다.</p>"

    def _convert_hwpx(self, file_path: str) -> str:
        """HWPX (Office Open XML 기반) 변환"""
        html_parts = []

        try:
            with zipfile.ZipFile(file_path) as zf:
                # Contents 디렉토리에서 section 파일 찾기
                section_files = sorted([f for f in zf.namelist()
                                       if 'Contents/section' in f and f.endswith('.xml')])

                if not section_files:
                    # 다른 구조 시도
                    section_files = [f for f in zf.namelist() if f.endswith('.xml') and 'section' in f.lower()]

                for section_file in section_files:
                    xml_content = zf.read(section_file)
                    tree = ET.fromstring(xml_content)

                    # HWPX 네임스페이스
                    ns = {
                        'hp': 'http://www.hancom.co.kr/hwpml/2011/paragraph',
                        'ht': 'http://www.hancom.co.kr/hwpml/2011/tail',
                        'hc': 'http://www.hancom.co.kr/hwpml/2011/core'
                    }

                    # 단락 추출
                    paragraphs = tree.findall('.//{http://www.hancom.co.kr/hwpml/2011/paragraph}p')
                    if not paragraphs:
                        # 네임스페이스 없이 시도
                        paragraphs = tree.findall('.//p')

                    for para in paragraphs:
                        # 텍스트 런 추출
                        texts = []
                        for elem in para.iter():
                            if elem.text:
                                texts.append(elem.text)
                            if elem.tail:
                                texts.append(elem.tail)

                        text = ''.join(texts).strip()
                        if text:
                            html_parts.append(f'<p>{text}</p>')

                    # 테이블 추출
                    tables = tree.findall('.//{http://www.hancom.co.kr/hwpml/2011/paragraph}tbl')
                    for tbl in tables:
                        html_parts.append(self._convert_hwpx_table(tbl))

        except Exception as e:
            return f"<p>HWPX 파싱 오류: {str(e)}</p>"

        return '\n'.join(html_parts) if html_parts else "<p>내용을 추출할 수 없습니다.</p>"

    def _convert_hwpx_table(self, tbl) -> str:
        """HWPX 테이블 변환"""
        html = '<table class="hwp-table">'
        rows = tbl.findall('.//{http://www.hancom.co.kr/hwpml/2011/paragraph}tr')

        for row in rows:
            html += '<tr>'
            cells = row.findall('.//{http://www.hancom.co.kr/hwpml/2011/paragraph}tc')

            for cell in cells:
                # 텍스트 추출
                texts = []
                for elem in cell.iter():
                    if elem.text:
                        texts.append(elem.text)
                text = ''.join(texts).strip()
                html += f'<td>{text}</td>'

            html += '</tr>'

        html += '</table>'
        return html

    # ============================================================
    # PDF 처리
    # ============================================================
    def _analyze_pdf(self, file_path: str) -> Tuple[bool, float]:
        """
        PDF 타입 분석

        Returns:
            (is_digital, text_ratio): 디지털 PDF 여부와 텍스트 비율
        """
        try:
            import pdfplumber
        except ImportError:
            # pdfplumber 없으면 무조건 Upstage 사용
            return False, 0.0

        try:
            with pdfplumber.open(file_path) as pdf:
                if not pdf.pages:
                    return False, 0.0

                # 처음 3페이지만 샘플링
                sample_pages = pdf.pages[:min(3, len(pdf.pages))]
                total_text_len = 0
                total_chars = 0

                for page in sample_pages:
                    text = page.extract_text() or ""
                    total_text_len += len(text.strip())

                    # 예상 문자 수 (페이지 크기 기반 추정)
                    width, height = page.width, page.height
                    expected_chars = (width * height) / 100  # 대략적인 추정
                    total_chars += expected_chars

                text_ratio = total_text_len / max(total_chars, 1)

                # 텍스트가 충분히 있으면 디지털 PDF
                is_digital = total_text_len > 100 and text_ratio > 0.3
                return is_digital, text_ratio

        except Exception:
            return False, 0.0

    def _convert_digital_pdf(self, file_path: str) -> str:
        """디지털 PDF를 HTML로 변환 (로컬 처리)"""
        try:
            import pdfplumber
        except ImportError:
            return "<p>pdfplumber 라이브러리가 필요합니다.</p>"

        html_parts = ['<div class="pdf-document">']

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                html_parts.append(f'<div class="pdf-page" data-page="{page_num}">')

                # 테이블 먼저 추출
                tables = page.extract_tables()
                table_bboxes = []

                for table in tables:
                    if table:
                        # 테이블 테두리 감지
                        has_lines = self._detect_table_lines(page, table)
                        border_class = 'bordered' if has_lines else 'borderless'

                        html_parts.append(f'<table class="pdf-table {border_class}">')
                        for row in table:
                            html_parts.append('<tr>')
                            for cell in row:
                                cell_text = self._clean_text(cell) if cell else ''
                                html_parts.append(f'<td>{cell_text}</td>')
                            html_parts.append('</tr>')
                        html_parts.append('</table>')

                # 일반 텍스트 추출
                text = page.extract_text()
                if text:
                    # 테이블에 포함된 텍스트 제거 후 단락으로 분리
                    paragraphs = text.split('\n\n')
                    for para in paragraphs:
                        para = para.strip()
                        if para:
                            html_parts.append(f'<p>{self._clean_text(para)}</p>')

                html_parts.append('</div>')
                html_parts.append('<hr class="page-break"/>')

        html_parts.append('</div>')
        return '\n'.join(html_parts)

    def _detect_table_lines(self, page, table) -> bool:
        """테이블의 선 유무 감지"""
        try:
            lines = page.lines or []
            rects = page.rects or []

            # 선이나 사각형이 있으면 테두리 있음
            return len(lines) > 0 or len(rects) > 0
        except:
            return True  # 기본값

    def _clean_text(self, text: str) -> str:
        """텍스트 정리"""
        if not text:
            return ""
        # HTML 특수문자 이스케이프
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        # 연속 공백 정리
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _convert_image_pdf_upstage(self, file_path: str) -> str:
        """
        이미지 PDF를 Upstage API로 변환

        Upstage Document Parse API 제한사항:
        - 최대 10페이지/요청 (대용량은 분할 필수)
        - 파일 크기 제한 (약 50MB)
        - Rate Limit 존재 (연속 요청 시 429 에러)

        크레딧 시스템: 1페이지당 55원 (부가세 포함)
        """
        import time as time_module
        import tempfile
        import shutil

        if not self.api_key:
            return "<p>Upstage API 키가 필요합니다. 이미지 PDF는 OCR이 필요합니다.</p>"

        # 파일 크기 확인
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > self.MAX_FILE_SIZE_MB:
            return f"<p>파일 크기가 너무 큽니다: {file_size_mb:.1f}MB (최대 {self.MAX_FILE_SIZE_MB}MB)</p>"

        try:
            from PyPDF2 import PdfReader, PdfWriter
        except ImportError:
            # 분할 없이 전체 업로드 (페이지 수 확인 불가)
            return self._call_upstage_api(file_path)

        # PDF 페이지 수 확인
        try:
            reader = PdfReader(file_path)
            total_pages = len(reader.pages)
        except Exception as e:
            return f"<p>PDF 파일을 읽을 수 없습니다: {str(e)}</p>"

        if total_pages == 0:
            return "<p>빈 PDF 파일입니다.</p>"

        # 크레딧 확인 (크레딧 시스템이 있는 경우)
        if self.credit_manager:
            has_credits, required, msg = self.credit_manager.check_credits(total_pages)
            if not has_credits:
                raise ValueError(f"크레딧 부족: {total_pages}페이지 변환에 {required:,}원 필요. 현재 잔액을 확인해주세요.")

        filename = os.path.basename(file_path)

        # 10페이지 이하는 그대로 처리
        if total_pages <= self.MAX_PDF_PAGES_PER_REQUEST:
            result = self._call_upstage_api(file_path)

            # 성공 시 크레딧 차감
            if self.credit_manager and result and not result.startswith("<p>"):
                self.credit_manager.deduct_credits(total_pages, filename)

            return result

        # === 대용량 파일 분할 처리 (10페이지씩) ===
        html_parts = []
        pages_processed = 0
        temp_files = []  # 임시 파일 목록 관리

        # 시스템 임시 디렉토리 사용 (Windows 호환성)
        temp_dir = tempfile.mkdtemp(prefix='lawpro_pdf_')

        try:
            total_chunks = (total_pages + self.MAX_PDF_PAGES_PER_REQUEST - 1) // self.MAX_PDF_PAGES_PER_REQUEST

            for chunk_idx, start_page in enumerate(range(0, total_pages, self.MAX_PDF_PAGES_PER_REQUEST)):
                end_page = min(start_page + self.MAX_PDF_PAGES_PER_REQUEST, total_pages)
                chunk_pages = end_page - start_page

                # 부분 PDF 생성
                writer = PdfWriter()
                for i in range(start_page, end_page):
                    writer.add_page(reader.pages[i])

                temp_path = os.path.join(temp_dir, f'chunk_{chunk_idx:03d}.pdf')
                with open(temp_path, 'wb') as f:
                    writer.write(f)
                temp_files.append(temp_path)

                # API 호출 (429 발생 시 내부에서 충분히 대기 후 재시도)
                try:
                    part_html = self._call_upstage_api(temp_path)

                    # 에러 응답 체크
                    if part_html.startswith("<p>API") or part_html.startswith("<p>오류"):
                        html_parts.append(f'<!-- Pages {start_page + 1}-{end_page}: ERROR -->')
                        html_parts.append(part_html)
                    else:
                        html_parts.append(f'<!-- Pages {start_page + 1}-{end_page} -->')
                        html_parts.append(part_html)
                        pages_processed += chunk_pages

                except Exception as e:
                    html_parts.append(f'<!-- Pages {start_page + 1}-{end_page}: FAILED -->')
                    html_parts.append(f'<p>페이지 {start_page + 1}-{end_page} 변환 실패: {str(e)}</p>')

            # 성공 시 크레딧 차감 (처리된 페이지만)
            if self.credit_manager and pages_processed > 0:
                self.credit_manager.deduct_credits(pages_processed, filename)

        finally:
            # 임시 파일 정리 (지연 삭제로 Windows 파일 잠금 문제 해결)
            time_module.sleep(0.5)  # 파일 핸들 해제 대기
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception:
                pass  # 정리 실패해도 계속 진행

        if not html_parts:
            return "<p>PDF 변환에 실패했습니다.</p>"

        return '\n'.join(html_parts)

    def _call_upstage_api(self, file_path: str) -> str:
        """
        Upstage Document Parse API 호출

        적응형 Rate Limit 시스템:
        - 요청 전: 쿨다운 상태 확인 → 잔여 시간만 대기
        - 성공 시: 쿨다운 리셋, 요청 빈도 기록 (학습 데이터)
        - 429 시: 쿨다운 설정, 요청 빈도 분석 및 Rate Limit 재학습

        쿨다운 시스템:
        - 429 발생 시 쿨다운 설정 (10초→30초→1분→2분→3분)
        - 다음 요청 시 잔여 시간만 대기 (블로킹 최소화)
        - 사용자가 프로세스 취소하면 자동 중단
        - 대기 중 아무 요청 없으면 다음 요청 시점에 시도

        에러 핸들링:
        - 429 Too Many Requests: 쿨다운 설정 후 재시도 (성공까지 반복)
        - 400 Bad Request: 파일 형식/크기 문제 (재시도 안함)
        - 500 Server Error: 대기 후 재시도 (최대 3회)
        """
        import sys
        import time as time_module

        filename = os.path.basename(file_path)

        headers = {
            "Authorization": f"Bearer {self.api_key}"
        }

        data = {
            "ocr": "force",
            "output_formats": '["html"]',
            "model": "document-parse",
            "coordinates": "false"
        }

        other_retry_count = 0
        max_other_retries = 3

        while True:
            try:
                # === 쿨다운 체크 (잔여 시간만 대기) ===
                if self.rate_limiter:
                    is_cooldown, remaining = self.rate_limiter.check_cooldown()
                    if is_cooldown and remaining > 0:
                        # 잔여 시간만 대기 (전체 대기 시간 아님)
                        log_msg = self.rate_limiter.get_cooldown_wait_log(remaining, filename)
                        print(log_msg, file=sys.stderr, flush=True)
                        time_module.sleep(remaining)

                    # Rate Limit 체크 (학습된 한도 기반 throttling)
                    should_wait, wait_sec = self.rate_limiter.should_wait()
                    if should_wait:
                        time_module.sleep(wait_sec)

                    # 요청 기록
                    self.rate_limiter.record_request()

                with open(file_path, "rb") as f:
                    files = {"document": (filename, f)}

                    response = requests.post(
                        self.UPSTAGE_API_URL,
                        headers=headers,
                        data=data,
                        files=files,
                        timeout=300
                    )

                # === 429 Too Many Requests ===
                if response.status_code == 429:
                    if self.rate_limiter:
                        # Rate Limit 분석 및 학습
                        analysis = self.rate_limiter.record_429_error()
                        log_msg = self.rate_limiter.get_429_analysis_log(analysis)
                        print(log_msg, file=sys.stderr, flush=True)

                        # 쿨다운 설정 (블로킹 없이 상태만 기록)
                        wait_time, retry_count = self.rate_limiter.set_cooldown(filename)
                        log_msg = self.rate_limiter.get_cooldown_log(wait_time, retry_count, filename)
                        print(log_msg, file=sys.stderr, flush=True)

                        # 쿨다운 시간만큼 대기 후 재시도
                        # (사용자가 취소하면 여기서 프로세스 종료됨)
                        time_module.sleep(wait_time)
                        continue
                    else:
                        # Rate Limiter 없으면 고정 대기
                        time_module.sleep(30)
                        continue

                # === 400 Bad Request ===
                if response.status_code == 400:
                    try:
                        error_detail = response.json()
                        error_msg = error_detail.get('message', error_detail.get('error', str(error_detail)))
                    except:
                        error_msg = response.text[:200]
                    return f"<p>API 요청 오류 (400): {error_msg}</p>"

                # === 500 Server Error ===
                if response.status_code >= 500:
                    other_retry_count += 1
                    if other_retry_count <= max_other_retries:
                        time_module.sleep(other_retry_count * 10)
                        continue
                    return f"<p>API 서버 오류 ({response.status_code}): 재시도 {max_other_retries}회 실패</p>"

                response.raise_for_status()
                result = response.json()

                # === 성공 ===
                if self.rate_limiter:
                    # 쿨다운 리셋 (성공했으므로)
                    self.rate_limiter.reset_cooldown()
                    # Rate Limit 학습 데이터 기록
                    self.rate_limiter.record_success()

                # HTML 추출
                if 'content' in result and 'html' in result['content']:
                    return result['content']['html']
                elif 'html' in result:
                    return result['html']
                elif 'text' in result:
                    return f"<p>{result['text']}</p>"
                else:
                    return f"<p>API 응답 형식 오류: {str(result)[:500]}</p>"

            except requests.exceptions.Timeout:
                other_retry_count += 1
                if other_retry_count <= max_other_retries:
                    time_module.sleep(10)
                    continue
                return "<p>API 요청 타임아웃 (5분 초과)</p>"

            except requests.exceptions.ConnectionError:
                other_retry_count += 1
                if other_retry_count <= max_other_retries:
                    time_module.sleep(5)
                    continue
                return "<p>API 연결 오류: 네트워크 확인 필요</p>"

            except requests.exceptions.HTTPError as e:
                other_retry_count += 1
                if other_retry_count <= max_other_retries:
                    time_module.sleep(5)
                    continue
                return f"<p>API HTTP 오류: {str(e)}</p>"

            except Exception as e:
                other_retry_count += 1
                if other_retry_count <= max_other_retries:
                    time_module.sleep(3)
                    continue
                return f"<p>API 오류: {str(e)}</p>"

    # ============================================================
    # HTML 저장
    # ============================================================
    def _save_html(self, save_path: str, content: str, original_filename: str):
        """완성된 HTML 저장"""
        html_template = f'''<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <meta name="generator" content="LawPro Fast Converter">
    <meta name="source-file" content="{original_filename}">
    <title>{original_filename}</title>
    <style>
        * {{
            box-sizing: border-box;
        }}

        body {{
            font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', sans-serif;
            line-height: 1.6;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
            background: #fff;
            color: #333;
        }}

        /* 테이블 공통 */
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 1em 0;
        }}

        td, th {{
            padding: 8px 12px;
            vertical-align: top;
        }}

        /* 테두리 있는 테이블 */
        table.bordered td,
        table.bordered th,
        table.excel-table td,
        table.excel-table th {{
            border: 1px solid #333;
        }}

        /* 테두리 없는 테이블 */
        table.borderless td,
        table.borderless th {{
            border: none;
        }}

        /* PDF 테이블 기본값 */
        table.pdf-table td {{
            border: 1px solid #ccc;
        }}

        table.pdf-table.borderless td {{
            border: none;
        }}

        /* 헤더 */
        h1, h2, h3 {{
            margin-top: 1.5em;
            margin-bottom: 0.5em;
            color: #222;
        }}

        h1 {{ font-size: 1.8em; border-bottom: 2px solid #333; padding-bottom: 0.3em; }}
        h2 {{ font-size: 1.4em; }}
        h3 {{ font-size: 1.2em; }}

        /* 단락 */
        p {{
            margin: 0.5em 0;
        }}

        /* 시트/슬라이드 구분 */
        .sheet, .slide {{
            margin: 2em 0;
            padding: 1em;
            background: #fafafa;
            border-radius: 4px;
        }}

        .sheet-title, .slide-number {{
            color: #666;
            font-size: 0.9em;
            margin-bottom: 1em;
        }}

        /* 페이지 구분선 */
        hr.page-break, hr.slide-divider {{
            border: none;
            border-top: 2px dashed #ccc;
            margin: 2em 0;
        }}

        /* 프린트 스타일 */
        @media print {{
            body {{
                max-width: none;
                padding: 0;
            }}
            hr.page-break {{
                page-break-after: always;
            }}
        }}
    </style>
</head>
<body>
    <header>
        <small style="color:#999;">Source: {original_filename}</small>
    </header>
    <main>
{content}
    </main>
    <footer style="margin-top:3em; padding-top:1em; border-top:1px solid #eee; color:#999; font-size:0.8em;">
        Converted by LawPro Fast Converter
    </footer>
</body>
</html>'''

        with open(save_path, 'w', encoding='utf-8') as f:
            f.write(html_template)
